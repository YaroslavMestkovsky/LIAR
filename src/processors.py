import logging
import re
import uuid

import PyPDF2
import openpyxl

from abc import ABC, abstractmethod
from pathlib import Path
from typing import List, Dict, Any, Optional
from docx import Document
from pptx import Presentation
from bs4 import BeautifulSoup

from qdrant_client import QdrantClient
from qdrant_client.models import PointStruct
from sentence_transformers import SentenceTransformer

from tqdm import tqdm

from src.models import init_embedding_model
from src.project_dataclasses import IndexingServiceConfig, BaseConfig


class BaseProcessor(ABC):
    """Базовый процессор."""

    def __init__(self, chunk_size, qdrant_client, config, base_config):
        self.qdrant: QdrantClient = qdrant_client
        self.base_config: BaseConfig = base_config
        self.config: IndexingServiceConfig = config
        self.logger: logging.Logger

        self.buffer: List[PointStruct] = []
        self.processed_count: int = 0
        self.chunk_size: int = chunk_size
        self.batch_size = self.config.batch_size

        self._setup_logger()

    def finalize(self):
        """Финализация обработки - отправка оставшихся данных."""

        if self.buffer:
            self._flush_batch()

        self.logger.debug(f"Обработка завершена. Всего обработано: {self.processed_count}")

    @abstractmethod
    def process_file(self, suffix: str, file_path: Path) -> bool:
        """Обработка одного файла."""

    @abstractmethod
    def process_directory(self, directory_path: Path) -> int:
        """Обработка директории."""

    def get_stats(self) -> Dict[str, Any]:
        """Получение статистики обработки."""

        return {
            "processed_count": self.processed_count,
            "buffer_size": len(self.buffer),
            "batch_size": self.config.batch_size,
        }

    @abstractmethod
    def _setup_logger(self) -> None:
        """Инициализация логгера."""

    def _flush_batch(self):
        """Отправка накопленного батча в Qdrant."""

        try:
            self.qdrant.upsert(
                collection_name=self.config.default_collection,
                points=self.buffer,
            )
            self.processed_count += len(self.buffer)
            self.logger.debug(f"Обработано {len(self.buffer)} элементов. Всего: {self.processed_count}")

        except Exception as e:
            self.logger.error(f"Ошибка при отправке батча: {e}", exc_info=True)
            raise
        finally:
            self.buffer = []

    def _create_point(
            self,
            vector: List[float],
            payload: Dict[str, Any],
            point_id: Optional[str] = None,
    ) -> PointStruct:
        """Создание точки для Qdrant."""

        if point_id is None:
            point_id = str(uuid.uuid4())

        return PointStruct(
            id=point_id,
            vector=vector,
            payload=payload,
        )

    def _add_to_buffer(self, point: PointStruct):
        """Добавление точки в буфер."""

        self.buffer.append(point)

        if len(self.buffer) >= self.batch_size:
            self._flush_batch()


class DocumentProcessor(BaseProcessor):
    """Процессор для обработки документов."""

    def __init__(self, chunk_size, qdrant_client, config, base_config):
        super().__init__(chunk_size, qdrant_client, config, base_config)
        self.embedding_model: SentenceTransformer = init_embedding_model(self.logger)

        # Поддерживаемые форматы
        self.document_formats = {
            '.pdf': self._extract_pdf_text,
            '.docx': self._extract_docx_text,
            '.doc': self._extract_doc_text,
            '.pptx': self._extract_pptx_text,
            '.ppt': self._extract_ppt_text,
            '.xlsx': self._extract_xlsx_text,
            '.xls': self._extract_xls_text,
            '.html': self._extract_html_text,
            '.htm': self._extract_html_text,
            '.txt': self._extract_txt_text,
            '.md': self._extract_md_text,
            '.rtf': self._extract_rtf_text,
        }

    def _setup_logger(self) -> None:
        """Настройка логирования."""

        log_config = self.base_config.logging
        log_params = {
            "encoding": "utf-8",
            "level": getattr(logging, log_config["level"]),
            "format": "[D.P.] %(asctime)s - %(levelname)s - %(message)s",
        }

        if log_config["log_in_file"]:
            log_params["filename"] = log_config["file"]

        logging.basicConfig(**log_params) # todo переписать на loguru и вынести в отдельный хелпер

        for name in ("httpx", "httpcore", "qdrant_client"):
            logging.getLogger(name).setLevel(logging.WARNING)

        self.logger = logging.getLogger(__name__)

    def _extract_pdf_text(self, file_path: Path) -> str:
        """Извлечение текста из PDF файла."""

        try:
            text_parts = []

            with open(file_path, 'rb') as file:
                pdf_reader = PyPDF2.PdfReader(file)

                for page_num, page in enumerate(pdf_reader.pages):
                    try:
                        page_text = page.extract_text()

                        if page_text.strip():
                            text_parts.append(page_text)

                    except Exception as e:
                        self.logger.warning(f"Ошибка при извлечении текста со страницы {page_num + 1}: {e}")
                        continue

            text = "\n".join(text_parts)

            self.logger.debug(f"Извлечено {len(text)} символов из PDF: {file_path}")
            return text

        except Exception as e:
            self.logger.error(f"Ошибка при извлечении текста из PDF {file_path}: {e}", exc_info=True)
            return ""

    def _extract_docx_text(self, file_path: Path) -> str:
        """Извлечение текста из DOCX файла."""

        try:
            doc = Document(file_path)
            text_parts = []

            # Извлечение текста из параграфов
            for paragraph in doc.paragraphs:
                if paragraph.text.strip():
                    text_parts.append(paragraph.text)

            # Извлечение текста из таблиц
            for table in doc.tables:
                for row in table.rows:
                    row_text = []

                    for cell in row.cells:
                        if cell.text.strip():
                            row_text.append(cell.text.strip())

                    if row_text:
                        text_parts.append(" | ".join(row_text))

            text = "\n".join(text_parts)

            self.logger.debug(f"Извлечено {len(text)} символов из DOCX: {file_path}")
            return text

        except Exception as e:
            self.logger.error(f"Ошибка при извлечении текста из DOCX {file_path}: {e}", exc_info=True)
            return ""

    def _extract_doc_text(self, file_path: Path) -> str:
        """Извлечение текста из DOC файла (старый формат)."""

        # Для старых DOC файлов нужна библиотека python-docx2txt или antiword
        # Пока возвращаем пустую строку
        self.logger.warning(f"Извлечение текста из DOC файлов не поддерживается: {file_path}")
        return ""

    def _extract_pptx_text(self, file_path: Path) -> str:
        """Извлечение текста из PPTX файла."""

        try:
            prs = Presentation(file_path)
            text_parts = []

            for slide_num, slide in enumerate(prs.slides):
                slide_text = []

                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_text.append(shape.text.strip())

                if slide_text:
                    text_parts.append(f"Слайд {slide_num + 1}: " + " | ".join(slide_text))

            text = "\n".join(text_parts)

            self.logger.debug(f"Извлечено {len(text)} символов из PPTX: {file_path}")
            return text

        except Exception as e:
            self.logger.error(f"Ошибка при извлечении текста из PPTX {file_path}: {e}", exc_info=True)
            return ""

    def _extract_ppt_text(self, file_path: Path) -> str:
        """Извлечение текста из PPT файла (старый формат)."""
        self.logger.warning(f"Извлечение текста из PPT файлов не поддерживается: {file_path}")
        return ""

    def _extract_xlsx_text(self, file_path: Path) -> str:
        """Извлечение текста из XLSX файла."""

        try:
            workbook = openpyxl.load_workbook(file_path, data_only=True)
            text_parts = []

            for sheet_name in workbook.sheetnames:
                sheet = workbook[sheet_name]
                sheet_text = []

                for row in sheet.iter_rows(values_only=True):
                    row_text = []

                    for cell in row:
                        if cell is not None and str(cell).strip():
                            row_text.append(str(cell).strip())

                    if row_text:
                        sheet_text.append(" | ".join(row_text))

                if sheet_text:
                    text_parts.append(f"Лист '{sheet_name}': " + "\n".join(sheet_text))

            text = "\n".join(text_parts)

            self.logger.debug(f"Извлечено {len(text)} символов из XLSX: {file_path}")
            return text

        except Exception as e:
            self.logger.error(f"Ошибка при извлечении текста из XLSX {file_path}: {e}", exc_info=True)
            return ""

    def _extract_xls_text(self, file_path: Path) -> str:
        """Извлечение текста из XLS файла (старый формат)."""
        self.logger.warning(f"Извлечение текста из XLS файлов не поддерживается: {file_path}")
        return ""

    def _extract_html_text(self, file_path: Path) -> str:
        """Извлечение текста из HTML файла."""

        try:
            with open(file_path, 'r', encoding='utf-8') as file:
                content = file.read()

            soup = BeautifulSoup(content, 'html.parser')

            # Удаление скриптов и стилей
            for script in soup(["script", "style"]):
                script.decompose()

            # Извлечение текста
            text = soup.get_text()

            # Очистка текста
            lines = (line.strip() for line in text.splitlines())
            chunks = (phrase.strip() for line in lines for phrase in line.split("  "))
            text = ' '.join(chunk for chunk in chunks if chunk)

            self.logger.debug(f"Извлечено {len(text)} символов из HTML: {file_path}")
            return text

        except Exception as e:
            self.logger.error(f"Ошибка при извлечении текста из HTML {file_path}: {e}", exc_info=True)
            return ""

    def _extract_txt_text(self, file_path: Path) -> str:
        """Извлечение текста из TXT файла."""

        try:
            # Попытка определить кодировку
            encodings = ['utf-8', 'cp1251', 'latin-1', 'iso-8859-1']

            for encoding in encodings:
                try:
                    with open(file_path, 'r', encoding=encoding) as file:
                        text = file.read()

                    self.logger.debug(f"Извлечено {len(text)} символов из TXT: {file_path}")
                    return text
                except UnicodeDecodeError:
                    continue

            self.logger.error(f"Не удалось определить кодировку файла: {file_path}")
            return ""

        except Exception as e:
            self.logger.error(f"Ошибка при извлечении текста из TXT {file_path}: {e}", exc_info=True)
            return ""

    def _extract_md_text(self, file_path: Path) -> str:
        """Извлечение текста из Markdown файла."""

        try:
            with open(file_path, 'r', encoding='utf-8') as file:
                text = file.read()

            # Удаление markdown разметки
            text = re.sub(r'#{1,6}\s+', '', text)  # Заголовки
            text = re.sub(r'\*\*(.*?)\*\*', r'\1', text)  # Жирный текст
            text = re.sub(r'\*(.*?)\*', r'\1', text)  # Курсив
            text = re.sub(r'`(.*?)`', r'\1', text)  # Код
            text = re.sub(r'\[(.*?)]\(.*?\)', r'\1', text)  # Ссылки
            text = re.sub(r'!\[.*?]\(.*?\)', '', text)  # Изображения

            self.logger.debug(f"Извлечено {len(text)} символов из MD: {file_path}")
            return text

        except Exception as e:
            self.logger.error(f"Ошибка при извлечении текста из MD {file_path}: {e}", exc_info=True)
            return ""

    def _extract_rtf_text(self, file_path: Path) -> str:
        """Извлечение текста из RTF файла."""

        try:
            with open(file_path, 'r', encoding='utf-8') as file:
                content = file.read()

            # Простое удаление RTF разметки
            text = re.sub(r'\\[a-z]+\d*\s?', '', content)
            text = re.sub(r'[{}]', '', text)
            text = re.sub(r'\s+', ' ', text)

            self.logger.debug(f"Извлечено {len(text)} символов из RTF: {file_path}")
            return text

        except Exception as e:
            self.logger.error(f"Ошибка при извлечении текста из RTF {file_path}: {e}", exc_info=True)
            return ""

    def process_file(self, suffix: str, file_path: Path):
        """Обработка одного документа."""

        try:
            self.logger.debug(f"Обработка документа: {file_path}")

            # Извлечение текста
            extractor = self.document_formats[suffix]
            text = extractor(file_path)

            if not text.strip():
                self.logger.warning(f"Не удалось извлечь текст из документа: {file_path}")
                return False

            # Разбиение на чанки
            chunks = self._chunk_text(text)

            if not chunks:
                self.logger.warning(f"Не удалось разбить текст на чанки: {file_path}")
                return False

            with tqdm(total=len(chunks), desc="Создание эмбедингов") as pbar:
                # Обработка каждого чанка
                for i, chunk in enumerate(chunks):
                    if not chunk.strip():
                        continue

                    # Создание эмбеддинга
                    embedding = self._create_embedding(chunk)

                    # Создание точки для Qdrant
                    payload = {
                        "file_path": str(file_path),
                        "file_type": "document",
                        "file_format": suffix,
                        "text": chunk,
                        "chunk_index": i,
                        "total_chunks": len(chunks),
                        "file_size": file_path.stat().st_size,
                    }

                    point = self._create_point(embedding, payload)
                    self._add_to_buffer(point)
                    pbar.update(1)

            self.logger.debug(f"Документ обработан: {file_path} ({len(chunks)} чанков)")
            return True

        except Exception as e:
            self.logger.error(f"Ошибка при обработке документа {file_path}: {e}", exc_info=True)
            return False

    def process_directory(self, directory_path: Path) -> int:
        pass

    def _chunk_text(self, text: str) -> List[str]:
        """Разбиение текста на чанки."""

        words = text.split()
        chunks = []

        current_chunk = []
        current_length = 0

        for word in words:
            if current_length + len(word) + 1 > self.chunk_size and current_chunk:
                chunks.append(" ".join(current_chunk))
                current_chunk = [word]
                current_length = len(word)
            else:
                current_chunk.append(word)
                current_length += len(word) + 1

        if current_chunk:
            chunks.append(" ".join(current_chunk))

        return chunks

    def _create_embedding(self, text: str) -> List[float]:
        """Создание эмбеддинга для текста."""

        try:
            if not text.strip():
                return [0.0] * 1024  # Пустой вектор

            embedding = self.embedding_model.encode(text, show_progress_bar=False)
            return embedding.tolist()

        except Exception as e:
            self.logger.error(f"Ошибка при создании эмбеддинга: {e}", exc_info=True)
            return [0.0] * 1024
